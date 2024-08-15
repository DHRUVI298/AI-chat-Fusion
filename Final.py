import streamlit as st
from dotenv import load_dotenv
import os
load_dotenv()
st.set_page_config(page_title="AI Chat Fusion",page_icon="💬")
st.title(":violet[AI Chat Fusion 💬]")

GROQ_API_KEY="EnterYourAPI"
st.sidebar.title(":violet[Navigation]")
# page = st.sidebar.selectbox(":blue[Choose a page]", ["ChatBoat", "Image","pdf","ppt","text","fuzzy","csv","translator","visualization","url","voiceassistance","video"])
page = st.sidebar.selectbox(":blue[Choose a page]", ["ChatBoat","pdf","voiceassistance","csv","text","video","translator","Image","ppt","visualization","url","Audio","fuzzy"])

#st.sidebar.selectbox("Choose a page", ["Page 1", "Page 2"])
theme = st.sidebar.selectbox(" :blue[Edit Theme :] ",("White","Black","Light Black","Light"))
#t2
if theme=="Black":
    page_bg_img = """
    <style>
        [data-testid = "stAppViewContainer" ]{
        background-color : #454545;
    }
    [data-testid= "stAppViewContainer"] div[role="button"],
    [data-testid= "stAppViewContainer"].st-emotion-cache-zt5igj e1nzilvr4{
    color: white !important;
    }
    </style>
    """
    st.markdown(page_bg_img,unsafe_allow_html=True)
if theme=="White":
    page_bg_img = """
    <style>
        [data-testid = "stAppViewContainer"]{
        background-color : rgba(240,240,240);
        }
    </style>
    """
    st.markdown(page_bg_img,unsafe_allow_html=True)
#t4
if theme=="Light Black":
    page_bg_img = """
    <style>
        [data-testid = "stAppViewContainer"]{
        background-color : rgba(100,100,100);
    }
    </style>
    """
    st.markdown(page_bg_img,unsafe_allow_html=True)
#t5
if theme=="Light":
    page_bg_img = """
    <style>
        [data-testid = "stAppViewContainer"]{
        background-color : white;
        color:black;
    }
    </style>
    """
    st.markdown(page_bg_img,unsafe_allow_html=True)
#data
     

#t1, t2,t3,t4,t5,t6  = st.tabs(['chatbot','imag-text','text-img','pdf','summary','text-com'])
if page == "ChatBoat":
    import os 
    #for the Api 
    # from groq import Groq
    #for gerate random variables
    import random
    # It offers a pre-built approach to manage the back-and-forth 
    # exchange between the user and the LLM model.
    from langchain.chains import conversation
    from langchain_groq import ChatGroq
    from langchain.chains.conversation.memory import ConversationBufferMemory
    from langchain.chains import ConversationChain
    from langchain.prompts import PromptTemplate
    from dotenv import load_dotenv
    import  http
    from groq import Groq
    import config    


    load_dotenv()
   
    groq_api_key = os.environ['GROQ_API_KEY']
    def main():

            
            
            st.subheader(":blue[ChatBoat App]")
            st.write("----")
            # st.sidebar.title("Select An LLm")
            # model = st.sidebar.selectbox(
            #     'choose a model',
            #     # ['Mixtral-8x7b-32768','llama2-70b-4096']
            #     (["Mixtral-8x7b-32768","llama2-70b-4096"])
                
            # )
            conversation_memory_length = st.sidebar.slider(':blue[Conversational Meomory Length:]',1,10,value=5)
            
            memory = ConversationBufferMemory(k=conversation_memory_length)
            
            user_question  = st.text_area(":black[Ask A Question...]")
            
            
            if 'chat_history' not in st.session_state:
                st.session_state.chat_history=[]
                #1
                st.session_state.selected_question = None
            else:
                for message in st.session_state.chat_history:
                    memory.save_context({'input':message['human']},{'output':message['AI']})
                    
            groq_chat = ChatGroq(        
            groq_api_key  = groq_api_key,
            model_name = "mixtral-8x7b-32768")
            
            conversation = ConversationChain(
                llm = groq_chat,
                memory = memory
                
            )
            if user_question:
                if st.session_state.selected_question is None:

                    response = conversation(user_question)
                    message = {'human':user_question,'AI':response['response']}

                    st.session_state.chat_history.append(message)

                    st.write("Chatbot : ",response['response'])
            
                else:
                    st.write(":green[select from chat history]")
            selected_question_index = st.selectbox('Select Question from History',[message['human'] for message in st.session_state.chat_history],index=None if st.session_state.selected_question is None else [message['human'] for message in st.session_state.chat_history].index(st.session_state.selected_question))


            if selected_question_index is not None :
                    st.session_state.selected_question = selected_question_index
                    selected_response = [message['AI'] for message in st.session_state.chat_history if message['human'] == selected_question_index][0]
                    st.write("Chatbot : ",selected_response)
        
            
            
    if __name__  == '__main__':
        main()
if page == "Image":
     import requests
     from PIL import Image
     import io
     import os 
     import base64
     API_URLL = "https://api-inference.huggingface.co/models/runwayml/stable-diffusion-v1-5"  
     API_URL = "https://api-inference.huggingface.co/models/Salesforce/blip-image-captioning-large"
     API_TOKEN = "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"
     
     def generate_image(prompt):
         headers = {"Authorization": f"Bearer {API_TOKEN}"}
         payload = {"inputs": prompt}
         response = requests.post(API_URLL, headers=headers, json=payload)
         return response.content
     
     def caption_image(image):
         headers = {"Authorization": f"Bearer {API_TOKEN}"}
         img_byte_arr = io.BytesIO()
         image.save(img_byte_arr, format="JPEG")
         img_byte_arr.seek(0)
         response = requests.post(API_URLL, headers=headers, data=img_byte_arr.getvalue())
         return response.json()
     
     st.subheader(":blue[Image processing Tools]")
     st.write("----")
     
     # User selection for task'
     tabs = st.tabs( ["Image Generation", "Image Captioning","Chat With Your Image"])
     
     with tabs[0]:
         
         # Text input for image generation
         st.subheader(":orange[Image Generation]")
         prompt = st.text_input("Enter a prompt for image generation:")
         if st.button("Generate Image"):
             image_bytes = generate_image(prompt)
     
             if image_bytes:
                 image = Image.open(io.BytesIO(image_bytes))
                 st.image(image, caption="Generated Image", use_column_width=True)
             else:
                 st.error("An error occurred while generating the image.")
     
     with tabs[1]:
         st.subheader(":orange[Image Captioning]")
         def query(image):
             headers = {"Authorization": f"Bearer {API_TOKEN}"}  # Use token if stored
             img_byte_arr = io.BytesIO()
             image.save(img_byte_arr, format="JPEG")
             img_byte_arr.seek(0)
             response = requests.post(API_URL, headers=headers, data=img_byte_arr.getvalue())
             return response.json()
     
     # Display file uploader and inference results
         uploaded_file = st.file_uploader("Upload an image", type=["jpg", "jpeg", "png"])
     
         if uploaded_file is not None:
             # Display the uploaded image
             image = Image.open(uploaded_file)
             
             st.image(image, caption="Uploaded Image", use_column_width=True)
     
             # Perform inference if user clicks the button
             if st.button("Run Inference"):
                 # Send image to model for inference
                 response = query(image)
     
                 # Display inference output
                 st.write("Inference Output:")
                 st.json(response)
     
     with tabs[2]:
         st.subheader(":orange[Chat With Your Image]")
         API_URLLL = "https://api-inference.huggingface.co/models/impira/layoutlm-document-qa"
         # Consider using environment variables for a slightly more secure token approach
         # (though st.secrets is still strongly recommended)
         API_TOKEN = "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"  # Replace with your actual token
         headers = {"Authorization": f"Bearer {API_TOKEN}"}
     
     
         def query(image_bytes, question):
             try:
             
                 image_b64 = base64.b64encode(image_bytes).decode()
                 payload = {"image": image_b64, "question": question}
                 response = requests.post(API_URLLL, headers=headers, json=payload)
                 return response.json()
             except Exception as e:
                 st.error(f"An unexpected error occurred during the API call: {e}")
     
     
         #st.subheader(":blue[Chat With Img]")
     
         image_file = st.file_uploader("Upload an image", type=["png", "jpg", "jpeg"])
         question = st.text_input(":black[Ask a question about the image:]")
     
         if st.button("Ask Question") and image_file is not None:
             with st.spinner("Processing image and question..."):
                 try:
                     image = Image.open(image_file)
                     st.image(image, caption=":orange[Uploaded Image]", use_column_width=True)
         
                     image_bytes = io.BytesIO()
                     image.save(image_bytes, format=image.format)
                     image_bytes = image_bytes.getvalue()
         
                     output = query(image_bytes, question)
            
                     # Check response type and handle accordingly
                     print(output)
                     st.success("Answer:")
                     st.write(output)
         
                 except Exception as e:
                     st.error(f"An error occurred: {e}")
         elif image_file is None:
             st.error("Please upload an image before asking a question.")
     
if page=="pdf":
       from dotenv import load_dotenv
       import streamlit as st
       import requests
       import streamlit as st
       import google.generativeai as genai  # pip install google-generativeai
       import re
       from langchain.text_splitter import CharacterTextSplitter
       from langchain.embeddings.huggingface import HuggingFaceEmbeddings
       from langchain.vectorstores import FAISS #facebook AI similarity search
       from langchain.chains.question_answering import load_qa_chain
       from langchain import HuggingFaceHub
       from fpdf import FPDF  # pip install FPDF
       from PyPDF2 import PdfReader # pip install PyPDF2
       from io import BytesIO
       import requests
       from PIL import Image
       import io       


       st.subheader(":blue[PDF processing Tools]")
       st.write("----")

       genai.configure(api_key="AIzaSyCuLcdytEEjNW0IAPja4UfrsCJyOoMY-JA")
       # User selection for task'
       tabs = st.tabs( ["PDF Generation", "Search In PDF","Chat With PDF","Join Pdf","Edit pdf"])
       
       with tabs[0]:
           def get_gemini_response(prompt):
               model = genai.GenerativeModel('gemini-pro')
               response = model.generate_content(prompt)
               return response
       
       
           def refine_subtopics(sub_topics, sub_titles):
               for sub_topic in sub_topics:
                   sub_titles.append(sub_topic[3:].replace('"', ""))
               return sub_titles
       
       
           def content_generation(sub_titles):
               content = []
               for i in sub_titles:
                   prompt = f"generate content of {i} for document on 2 bullet point only each of point 20 tokens"
                   model = genai.GenerativeModel('gemini-pro')
                   response = model.generate_content(prompt)
                   content.append(response.text)
               return content


           def clean_text(text):
               cleaned_text = re.sub('\s+', ' ', text).strip()
               cleaned_text = re.sub(r'[-]\s|\d+\.\s*', '', cleaned_text)
               cleaned_text = re.sub(r'\s*:\s*', ': ', cleaned_text)
               cleaned_text = re.sub(r'\s*-\s*', ' - ', cleaned_text)
               return cleaned_text


           def split_sentences(text):
               sentences = re.split(r'(?<=\.)\s+', text)
               sentences = [sentence.capitalize() for sentence in sentences]
               return sentences
       

           def refine_fine_content(content):
               final_content = []
               for i in content:
                   cleaned_text = clean_text(i)
                   sentences = split_sentences(cleaned_text)
                   final_content.append(sentences)
               return final_content


           def create_pdf(topic, sub_titles, final_content):
               pdf = FPDF()
               pdf.add_page()
               pdf.set_font("Arial", size=12)
       
               # Add Title
               pdf.cell(200, 10, txt=topic, ln=1, align='C')
               pdf.set_font("Arial", size=10)
       
               # Add Subtopics and Content
               for i, sub_title in enumerate(sub_titles):
                   pdf.cell(0, 6, txt=f"{i+1}. {sub_title}", ln=1)
                   for point in final_content[i]:
                       pdf.write(5, point)
       
               pdf.output(f"E:\Python\document_{topic}.pdf")
               st.success(f"PDF document created successfully: document_{topic}.pdf")
       

           def download_button(file_path, topic):
               with open(file_path, "rb") as file:
                   pdf_content = file.read()
       
               st.download_button(
                   label="Download PDF",
                   data=pdf_content,
                   file_name=f"{topic}.pdf",
                   key="pdf_download_button"
               )
       
       
           st.subheader(":orange[PDF Genration]")
           topic = st.text_input("Enter topic : ", key="input")
           no_of_slides = st.text_input("Enter number of subtopics (for reference): ", key="slide")
       
           submit = st.button("Generate PDF")
       
           if submit:
               prompt = f"Generate a {no_of_slides} sub-titles only on the topic of {topic}"
               response = get_gemini_response(prompt)
               sub_topics = response.text.split("\n")
               sub_titles = refine_subtopics(sub_topics, [])
               content = content_generation(sub_titles)
               final_content = refine_fine_content(content)
       
               create_pdf(topic, sub_titles, final_content)
               download_button(f"E:\Python\document_{topic}.pdf", topic)

        

       with tabs[1]:      
           def search_pdf(file,query):
               #text = ""
               reader = PdfReader(file)
               found=False
               #num_pages = len(reader.pages)
               
               for page_num, page in enumerate(reader.pages):
                       #page = reader.pages[page_num]
                       text = page.extract_text()
       
                       paragraphs = text.split("\n")
       
                       current_title = None
                       current_content = ""
                       for paragraph in paragraphs:
                           
                           if query.lower() in paragraph.lower():
                               current_title = paragraph.strip()
                               current_content = ""
                           elif current_content is not None:
                               current_content += paragraph.strip() + "\n"
       
                       if current_title is not None and query.lower() in current_title.lower():
                               st.write("Answer Found on Page : ",page_num+1)
                               st.write("Title : ",current_title)
                               st.write("content : ",current_content)
                               #st.write(text)
                               found=True
                               break
               if not found:
                   st.write("Answer not Found in PDF")
                   
           def main():
               st.subheader(":orange[search With PDF]")
               upload_file = st.file_uploader(":orange[Upload a PDF : ]",type="pdf")
               if upload_file is not None:
                   st.write(":green[PDF uploaded successfully]")
                   query = st.text_input(":black[Enter Your Question : ]")
                   if st.button("Search"):
                       pdf_content = BytesIO(upload_file.read())
                       search_pdf(pdf_content,query)
       
           if __name__ == "__main__":
               main()
               
       with tabs[2]:
           def main():
               load_dotenv()
               #st.set_page_config(page_title="Ask your PDF")
               st.subheader(":orange[Ask Your PDF]")
       
               pdf = st.file_uploader(":black[Upload your pdf]",type="pdf")
           
               if pdf is not None:
                   pdf_reader = PdfReader(pdf)
                   text = ""
                   for page in pdf_reader.pages:
                       text += page.extract_text()
           
                   # spilit ito chuncks
                   text_splitter = CharacterTextSplitter(
                       separator="\n",
                       chunk_size=1000,
                       chunk_overlap=200,
                       length_function=len
                   )
                   chunks = text_splitter.split_text(text)
           
                   # create embedding
                   embeddings = HuggingFaceEmbeddings()
           
                   knowledge_base = FAISS.from_texts(chunks,embeddings)
           
                   user_question = st.text_input(":black[Ask Question about your PDF:]")
                   if user_question:
                       docs = knowledge_base.similarity_search(user_question)
                       llm = HuggingFaceHub(repo_id="google/flan-t5-large", model_kwargs={"temperature":5,
                                                                 "max_length":64})
                       chain = load_qa_chain(llm,chain_type="stuff")
                       response = chain.run(input_documents=docs,question=user_question)
           
                       st.write(response)
           
    
           if __name__ == '__main__':
               main()
    
            # st.write(chunks)
       with tabs[3]:
            import streamlit as st
            from PyPDF2 import PdfWriter
            from io import BytesIO

            st.title("PDF Merger")

            pdf_files = st.file_uploader("Upload PDF : ",accept_multiple_files=True,type="pdf")


            if st.button("Merge PDFs"):
                if pdf_files:
                    merger = PdfWriter()

                    for pdf_file in pdf_files:
                        merger.append(pdf_file)

                    output = BytesIO()
                    merger.write(output)

                    st.download_button(label="Download Merged PDF",data=output.getvalue(),file_name="merged.pdf",mime="application/pdf")
       with tabs[4]:
            import streamlit as st
            import fitz  # PyMuPDF library

            def edit_pdf_text(pdf_file, old_text, new_text):
                """Edits text in a PDF by replacing all occurrences of a specified text.

                Args:
                    pdf_file (streamlit.uploadedfile.UploadedFile): The uploaded PDF file.
                    old_text (str): The text to be replaced.
                    new_text (str): The new text to replace all occurrences of the specified text.

                Returns:
                    bytes: The edited PDF content in bytes format.
                """
                try:
                    # Open the PDF document
                    doc = fitz.open(stream=pdf_file.read())

                    # Iterate through all pages and search for the specified text
                    for page in doc:
                        # Search for text based on coordinates
                        blocks = page.search_for(old_text)
                        if blocks:
                            # Iterate through each block and draw a rectangle over the old text
                            for block in blocks:
                                x0, y0, x1, y1 = block
                                page.add_redact_annot((x0, y0, x1, y1))
                                page.apply_redactions()
                                page.insert_text(block[:2], new_text)  # Replace text at specified location

                    # Save the edited PDF
                    edited_pdf_data = doc.write()
                    return edited_pdf_data

                except Exception as e:
                    st.error(f"Error editing PDF: {e}")
                    return None

            def main():
                """Builds the Streamlit app for PDF text editing."""
                st.title("PDF Text Editor")
                uploaded_file = st.file_uploader("Upload PDF", type="pdf")

                # User input for text editing
                old_text = st.text_input("Enter text to be replaced")
                new_text = st.text_input("Enter new text")

                # Download edited PDF
                if uploaded_file is not None and old_text and new_text:
                    edited_pdf = edit_pdf_text(uploaded_file, old_text, new_text)
                    if edited_pdf:
                        st.download_button(label="Download Edited PDF", data=edited_pdf, file_name="edited.pdf")

            if __name__ == "__main__":
                main()

           
           
               
if page=="ppt":
      import google.generativeai as genai  # pip install google-generativeai
   # pip install python-pptx
      import re
      genai.configure(api_key="AIzaSyCuLcdytEEjNW0IAPja4UfrsCJyOoMY-JA")
      def get_gemini_response(prompt):
          model = genai.GenerativeModel('gemini-pro')
          response = model.generate_content(prompt)
          return response
      sub_titles = []
      def refine_subtopics(sub_topics, sub_titles):
          for sub_topic in sub_topics:
              sub_titles.append(sub_topic[3:].replace('"',""))
          return sub_titles

      content = []
      def content_generation(sub_titles):
          for i in sub_titles:
              prompt = f"generate contebt of {i} for presentation slide on 2 bullet point only each of point 20 tokens"
              model = genai.GenerativeModel('gemini-pro')
              response = model.generate_content(prompt)
              content.append(response.text)  
             # content.append(response)
              
          return content

      def clean_text(text):
          cleaned_text = re.sub('\s+', ' ', text).strip()
          cleaned_text = re.sub(r'[*-]\s*|\d+\.\s*', '', cleaned_text)
          cleaned_text = re.sub(r'\s*:\s*', ': ',cleaned_text)
          cleaned_text = re.sub(r'\s*-\s*', ' - ',cleaned_text)
          return cleaned_text

      def split_sentences(text):
          sentences = re.split(r'(?<=\.)\s+', text)
          sentences = [sentence.capitalize() for sentence in sentences]
          return sentences

      def replace_and_capitalize(text):
          def replace_and_capitalize_colon(match):
              return match.group(1) + match.group(2).capitalize() + match.group(3)
          
          result = re.sub(r'(:\s*)(.*?)(\s*:[^:]|$)',replace_and_capitalize_colon, text)
          return result

      final_content = []
      def refine_fine_content(content):
          for i in content:
              cleaned_text = clean_text(i)
              sentences = split_sentences(cleaned_text)
              final_content.append(sentences)
          print("final content ready..")
          print(final_content)
          return final_content

      from pptx import Presentation
      powerpoint = Presentation()

      def slide_maker(powerpoint, topic,sub_titles, final_content):
      
          title_slide_layout = powerpoint.slide_layouts[0]
          title_slide = powerpoint.slides.add_slide(title_slide_layout)
          title = title_slide.shapes.title
          title.text = topic
          #title.text_frame.paragraphs[0].font.size = Pt(32)
          title.text_frame.paragraphs[0].font.bold = True
          content = title_slide.placeholders[1]
          content.text = "Created by genai"
    

          for i in range(len(sub_titles)):
      
              bulletLayout = powerpoint.slide_layouts[1]
              secondSlide = powerpoint.slides.add_slide(bulletLayout)
      
              myShapes = secondSlide.shapes
              titleShape = myShapes.title
              bodyShape = myShapes.placeholders[1]
              titleShape.text = sub_titles[i]
           #   titleShape.text_frame.paragraphs[0].font.size = Pt(24)
              titleShape.text_frame.paragraphs[0].font.bold = True
              tFrame = bodyShape.text_frame
      
              print("Topic generate")
              print(topic)
      
              for point in final_content[i]:
                  point = re.sub(r':[^:]+', ':', point)
                  point = replace_and_capitalize(point)
                  p = tFrame.add_paragraph()
                  p.text = point
                  #p.font.size = Pt(18)
                  #p.space_after = Pt(10)
          return powerpoint

      def download_button(file_path,topic):
          with open(file_path, "rb") as file:
              ppt_content = file.read()
      
          st.download_button(
              label="Download PPT",
              data = ppt_content,
              file_name=f"{topic}.pptx",
              key = "ppt_download_button"
          )

      st.subheader(":blue[Generate your ppt]")
      st.write("----")
      topic = st.text_input("Enter topic : ",key="input")
      no_of_slides = st.text_input("Enter number of slides : ",key="slide")

      submit = st.button("Generate PPT")

      if submit:
          prompt = f"Generate a {no_of_slides} sub-titles only on the topic of {topic}"
          response = get_gemini_response(prompt)
          print("Topic Generate")
          sub_topics = response.text.split("\n")
          sub_titles = refine_subtopics(sub_topics, sub_titles)
          print("Sub Titles")
          content = content_generation(sub_titles)
          print("content generated")
    
          #final_content = refine_fine_content(content) change 2
          final_content = refine_fine_content(content)
          print("fina; content ready")
          powerpoint = slide_maker(powerpoint,topic, sub_titles, final_content)
          powerpoint.save(f"E:\Final_Project_Sem8\PowerPoint\{topic}.pptx")
          st.text("PPT Created Successfully")
          download_button(f"E:\Final_Project_Sem8\PowerPoint\{topic}.pptx",topic)
          print("DONE finally ")


if page=="text":
     import os 
     import base64
     from langchain.text_splitter import CharacterTextSplitter
     from langchain.embeddings.huggingface import HuggingFaceEmbeddings
     from langchain.vectorstores import FAISS #facebook AI similarity search
     from langchain.chains.question_answering import load_qa_chain
     from langchain import HuggingFaceHub

#for the Api 
     from groq import Groq
#for gerate random variables
     import random
# It offers a pre-built approach to manage the back-and-forth 
# exchange between the user and the LLM model.
     from langchain.chains import conversation
     from langchain.chains.conversation.memory import ConversationBufferMemory
     from langchain_groq import ChatGroq
     from langchain.chains import ConversationChain
     from langchain.prompts import PromptTemplate
     from dotenv import load_dotenv
     from groq import Groq
     import config
     import streamlit as st
     from PyPDF2 import PdfReader # pip install PyPDF2
     from io import BytesIO
     import streamlit as st
     import requests
     from PIL import Image
     import io
     import streamlit as st
     import requests
     import json
     
     st.subheader(":blue[Text Automation Tools]")
     st.write("----")
     tabs = st.tabs( ["Text sumary","Auto Fill"])
     with tabs[0]:
        
         API_URL = "https://api-inference.huggingface.co/models/facebook/bart-large-cnn"
         API_TOKEN = "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"
     #Text Summarization with BART Model
         st.subheader(":orange[Text summary]")

         # Function to send text to the model and get the summarization response
         def query_bart(text):
             headers = {"Authorization": f"Bearer {API_TOKEN}"}
             payload = {"inputs": text}
             response = requests.post(API_URL, headers=headers, json=payload)
             return response.json()

         # Input field for the user to enter text
         text_input = st.text_area(":black[Enter text to summarize:]")

    # Perform summarization if user clicks the button
         if st.button("Summarize"):
             if text_input:
                 # Send user input text to the model for summarization
                 response = query_bart(text_input)
     
                 # Display the full response
                 st.write("Full Response:", response)
     
                 # Display the summarization result if available
                 if "summary_text" in response:
                     summary = response["summary_text"]
                     st.subheader("Summary:")
                     st.write(summary)
                 else:
                     st.error("summary from the response.")
             else:
                 st.warning("Please enter some text to summarize.")
     with tabs[1]:
             # Adjust API URL and headers if needed for a different model
             # f"https://api-inference.huggingface.co/models/{model_id}
         API_URL = "https://api-inference.huggingface.co/models/google/gemma-7b"
         API_TOKEN = "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"  # Replace with your actual token
         
         def query(prompt):
             headers = {"Authorization": f"Bearer {API_TOKEN}"}
             payload = {"inputs": prompt}
             try:
                 response = requests.post(API_URL, headers=headers, json=payload)
                 response.raise_for_status()  # Raise an exception for non-200 status codes
                 return response.json()
             except requests.exceptions.RequestException as e:
                 st.error(f"An error occurred while communicating with the API: {e}")
                 return None
        

         def main():
             #Gemma-7B Text Generation with Streamlit
             st.subheader(":orange[Auto generated Text]")
     
             prompt_text = st.text_input(":black[Enter a prompt to continue the sentence:]")
             if prompt_text:
                 completed_text = query(prompt_text)
     
                 if completed_text:
                     st.success("**Generated Text:**")
                     if isinstance(completed_text, list):
                         st.write(completed_text[0])  # Access text from the first element of the list
                     else:
                         st.write(completed_text.get("generated_text", ""))  # Original logic for dictionaries
                 else:
                     st.warning("The model did not generate any text for this prompt.")
     
             # st.write("**Note:**")
             # st.write("- Ensure you have a valid API key from Hugging Face to use this app.")
             # st.write("- This model might perform better with specific prompts or sentence starters.")
     
         if __name__ == "__main__":
             main()

if page=="fuzzy":
      st.subheader(":blue[Fuzzy Logic Related Tools]")
      st.write("----")
      import cv2
      import numpy as np 
      import skfuzzy as fuzz      
      tabs = st.tabs( ["Image", "text"])
     
      with tabs[0]:
              def brightness_classification(image):
               
                    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
                    # Calculate mean brightness   ----------> FUZZY LOGIC
                    mean_brightness = np.mean(gray_image)
           
                    # Define fuzzy membership functions for brightness
                    brightness = np.arange(0, 256, 1)
                    brightness_low = fuzz.trimf(brightness, [0, 0, 127])
                    brightness_high = fuzz.trimf(brightness, [128, 255, 255])
              
                    membership_low = fuzz.interp_membership(brightness, brightness_low, mean_brightness)
                    membership_high = fuzz.interp_membership(brightness, brightness_high, mean_brightness)
               
                    # Output classification result   --------------> DEFUZZY LOGIC
                    if membership_low > membership_high:
                        st.markdown(f"<h4 style='color:blue;'> Image contains Low Brightness</h4>",unsafe_allow_html=True)
           
                    else:
                        st.markdown(f"<h4 style='color:blue;'> Image contains High Brightness</h4>",unsafe_allow_html=True)
           
           
           
              def texture_classification(image):
              
                  gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
             
                  laplacian = cv2.Laplacian(gray_image, cv2.CV_64F)
             
                  mean_texture = np.mean(laplacian)    #-------------> FUZZY LOGIC
              
                  texture = np.arange(0, 256, 1)
                  texture_low = fuzz.trimf(texture, [0, 0, 127])
                  texture_high = fuzz.trimf(texture, [128, 255, 255])
              
                  membership_low = fuzz.interp_membership(texture, texture_low, mean_texture)
                  membership_high = fuzz.interp_membership(texture, texture_high, mean_texture)
           
                  # Output classification result    -------------> DEFUZZY LOGIC
                  if membership_low > membership_high:
                      st.markdown(f"<h4 style='color:blue;'> Image contains Low Texture</h4>",unsafe_allow_html=True)
              
                  else:
                      st.markdown(f"<h4 style='color:blue;'> Image contains High Texture</h4>",unsafe_allow_html=True)
           
     
     
              def main():
                   st.subheader(":orange[Image Classification through Fuzzy Logic]")
                   st.markdown("---")
           
                   uploaded_file = st.file_uploader("Upload an image", type=["png", "jpg", "jpeg"])
                   if uploaded_file is not None:
                       image = cv2.imdecode(np.frombuffer(uploaded_file.read(), np.uint8), 1)
                       st.markdown("---")
                       st.markdown("1. User can upload any format of image such as png, jpg, jpeg")
                       st.markdown("2. This classification implements both Fuzzification and defuzzification.")
           
                       st.image(image, caption="Uploaded Image", use_column_width=True)
                       brightness_result = brightness_classification(image)
                       texture_result = texture_classification(image)
                   
           
              if __name__ == "__main__":
                  main()

      with tabs[1]:
            from fuzzywuzzy import fuzz
            st.subheader(":orange[Text Flex with fuzzy logic]")
            st.markdown("---")
            st.markdown("Optimizing User Interaction with Fuzzy Text Matching")
            st.markdown("---")
            st.markdown("1. Fuzzy Ratio") 
            st.markdown("2. Fuzzy Partial Ratio") 
            st.markdown("3. Token Sort Ration")
            st.markdown("4. Token Set Ratio")
            st.markdown("5. Closet Match Ratio")
            st.markdown("---")
      
      
            text1 = st.text_input(":orange[Enter First Input Text:]")
            st.markdown("---")
            text2 = st.text_input(":orange[Enter Second Input Text:]")
            st.markdown("---")
            ratio = fuzz.ratio(text1.lower(),text2.lower())
            st.markdown(f"<h3 style='color:blue;'>Fuzzy Ration of inputs : {ratio}</h3>",unsafe_allow_html=True)
            ration = fuzz.partial_ratio(text1.lower(),text2.lower())
            st.markdown(f"<h3 style='color:blue;'>Fuzzy Partial Ration : {ration}</h3>",unsafe_allow_html=True)
            token_sort_ratio = fuzz.token_sort_ratio(text1, text2)
            st.markdown(f"<h3 style='color:blue;'>Token Sort Ratio : {token_sort_ratio}</h3>",unsafe_allow_html=True)
            token_set_ratio = fuzz.token_set_ratio(text1, text2)
            st.markdown(f"<h3 style='color:blue;'>Token Set Ratio : {token_set_ratio}</h3>",unsafe_allow_html=True)
            from fuzzywuzzy import process
     
if page=="csv":
     import pandas as pd

     def load_csv(file):
         try:
             df = pd.read_csv(file)
             return df
         except pd.errors.EmptyDataError:
             st.error("The uploaded CSV file is empty.")
             return None
         except Exception as e:
             st.error(f"Error occurred: {e}")
             return None
     
     # Function to search for keyword in CSV
     def search_csv(df, keyword):
         results = df.apply(lambda row: row.astype(str).str.contains(keyword, case=False)).any(axis=1)
         search_results = df[results]
         return search_results
     
     # Main function
     def main():
         st.subheader(":blue[CSV Viewer and Search]")
         
         st.markdown("Upload CSV File")
         uploaded_file = st.file_uploader("Upload CSV", type=['csv'])
         
         if uploaded_file is not None:
             st.success('File successfully uploaded.')
             df = load_csv(uploaded_file)
             
             if df is not None:
                 st.subheader("View CSV")
                 st.write(df)
                 
                 st.subheader("Search CSV")
                 user_input = st.text_input("Search for a keyword: ")
                 if st.button("Search"):
                     if user_input:
                         search_results = search_csv(df, user_input)
                         if not search_results.empty:
                             st.write("Search Results:")
                             st.write(search_results)
                         else:
                             st.write("No results found for the keyword.")
     
     if __name__ == "__main__":
         main()
if page=="translator":
      
     from translate import Translator

# Streamlit UI
     st.subheader(":blue[Text Translator]")

# Text input
     text = st.text_area("Enter text to translate", "")

# Language selection
     source_lang = st.selectbox("Select source language", options=[ 'en', 'es', 'fr', 'de', 'ja', 'kn', 'hi', 'mr', 'ta', 'te', 'gu', 'ko', 'ru', 'ar', 'it'])
     dest_lang = st.selectbox("Select destinaton language", options=['en', 'es', 'fr', 'de', 'ja', 'ko', 'hi', 'kn', 'mr', 'ta', 'te', 'gu', 'ru', 'ar', 'it'])

# Translate button
     if st.button("Translate"):
    # Create an instance of the Translator class
        translator = Translator(to_lang=dest_lang, from_lang=source_lang)

    # Translate text
        translated_text = translator.translate(text)
    
    # Display translated text
        st.write("Translated Text:")
        st.write(translated_text)


if page=="visualization":
     import pandas as pd

     from io import StringIO
     chart_type_map = {
       "bar": st.bar_chart,
       "line": st.line_chart,
       
       "area": st.area_chart,
       "scatter": st.scatter_chart,
       # Add more mappings for other chart types (e.g., histogram, boxplot)
     }

     import pandas as pd
     import streamlit as st
     from io import StringIO

# Function to generate the chart based on data, chart type, and selected columns
     def generate_chart(data, chart_type, selected_columns):
       if not selected_columns:
         st.error("Please select at least one column to visualize.")
         return
     
       try:
         chart_type_map[chart_type](data[selected_columns])  # Use selected columns
       except ValueError:
         st.error("Invalid data types for selected columns. Please try different columns or chart types.")

     # App layout
     st.subheader(":blue[CSV to Chart Visualization]")
     st.write("Upload your CSV file, select columns, and choose a chart type")

# File uploader
     uploaded_file = st.file_uploader("Choose a CSV file", type="csv")

     if uploaded_file is not None:
  # Read CSV data
       data = pd.read_csv(StringIO(uploaded_file.read().decode("utf-8")))
       st.success("CSV file uploaded successfully!")

  # Display a sample of the data
       st.dataframe(data.head(10))

  # Allow multi-select column selection
       selected_columns = st.multiselect("Select columns to visualize", options=data.columns, default=[])

  # Chart type selection
       chart_type = st.selectbox("Select chart type", list(chart_type_map.keys()))

  # Check for large dataset and offer sampling option
       if len(data) > 10000:
         if st.checkbox("Dataset is large. Show a random sample instead?"):
           data = data.sample(10000)
           st.warning("Displaying a random sample of 10,000 rows for performance.")
     
       # Generate chart
       if st.button("Generate Chart"):
         generate_chart(data, chart_type, selected_columns)

if page=="url":
       st.subheader(":blue[URL Related Tools]")
       import os 
       import base64
       from langchain.text_splitter import CharacterTextSplitter
       from langchain.embeddings.huggingface import HuggingFaceEmbeddings
       from langchain.vectorstores import FAISS #facebook AI similarity search
       from langchain.chains.question_answering import load_qa_chain
       from langchain import HuggingFaceHub
       from textblob import TextBlob
       import streamlit as st
       from urllib.request import urlopen
#for the Api 
       from groq import Groq
#for gerate random variables
       import random
# It offers a pre-built approach to manage the back-and-forth 
# exchange between the user and the LLM model.
       from langchain.chains import conversation
       from langchain.chains.conversation.memory import ConversationBufferMemory
       from langchain_groq import ChatGroq
       from langchain.chains import ConversationChain
       from langchain.prompts import PromptTemplate
       from dotenv import load_dotenv
       from groq import Groq
       import config
       import streamlit as st
       from PyPDF2 import PdfReader # pip install PyPDF2
       from io import BytesIO
       import streamlit as st
       import requests
       from PIL import Image
       import io
       import streamlit as st
       import requests
       import json
       import streamlit as st
       from bs4 import BeautifulSoup
       import requests
       from sumy.parsers.plaintext import PlaintextParser
       import sumy.nlp.tokenizers 
       from sumy.summarizers.lex_rank import LexRankSummarizer
       import nltk       


       nltk.download('punkt')
       tabs = st.tabs( ["weburl summary","Analyzes sentiments","Website Design Assistant"])
       with tabs[0]:
            def summarize_text(text, num_sentences=3):
                """Summarizes a given text using the LexRank summarizer from the sumy library.
            
                Args:
                    text: The text to be summarized (str).
                    num_sentences: The desired number of sentences in the summary (int, default=3).
            
                Returns:
                    A string containing the summarized text.
                """
                parser = PlaintextParser.from_string(text, sumy.nlp.tokenizers.Tokenizer("english"))
                summarizer = LexRankSummarizer()
            
                summary = summarizer(parser.document, sentences_count=num_sentences)
            
                return '\n'.join([str(sentence) for sentence in summary])  # Join sentences with line breaks
            
            def summarize_website(url):
                """
                Fetches website content, extracts text, and summarizes it using LexRank.
                """
                try:
                    response = requests.get(url)
                    response.raise_for_status()  # Raise an exception for non-200 status codes
            
                    soup = BeautifulSoup(response.content, 'html.parser')
                    text = soup.get_text(separator='\n')
            
                    summary = summarize_text(text)
                    return summary
            
                except requests.exceptions.RequestException as e:
                    return f"Error: Failed to access website ({e})"
            
            st.subheader(":orange[Chat with Website]")
            website_url = st.text_input("Enter website URL:")
            
            if website_url:
                summary = summarize_website(website_url)
                st.write(f"Summary of {website_url}:")
                st.write(summary)
                st.write("**Note:** This is a basic example and might not work for all websites.")  

       with tabs[1]:
            def analyze_sentiment(text):
              """
              Analyzes sentiment using TextBlob library.
              """
              analysis = TextBlob(text)
              if analysis.sentiment.polarity > 0:
                return "Positive"
              elif analysis.sentiment.polarity == 0:
                return "Neutral"
              else:
                return "Negative"
            
            def fetch_website_content(url):
              """
              Fetches the content of a website using urlopen.
              - Handles potential exceptions (e.g., invalid URL, network errors).
              """
              try:
                response = urlopen(url)
                return response.read().decode('utf-8')
              except Exception as e:
                st.error(f"Error fetching content: {e}")
                return None
            
            st.subheader(":orange[Website Sentiment Analysis (TextBlob)]")
            st.write("Enter a website URL or short text content")
            
            user_input = st.text_input("", key="user_input")  # Use key to avoid state conflicts
            
            if st.button("Analyze Sentiment"):
              if user_input:
                if user_input.startswith("http"):  # Check if it's a URL
                  website_content = fetch_website_content(user_input)
                  if website_content:
                    sentiment = analyze_sentiment(website_content)
                    st.write(f"Sentiment of website content: {sentiment}")
                  else:  # Handle content fetching errors
                    st.warning("Failed to fetch content. Please check the URL.")
                else:  # Analyze text directly
                  sentiment = analyze_sentiment(user_input)
                  st.write(f"Sentiment of text: {sentiment}")
              else:
                st.warning("Please enter a website URL or text content to analyze.")
                
       with tabs[2]:
            import streamlit as st
            from random import choices  # For generating sample content suggestions

        # Website information gathering
            st.title("Website Design Assistant")
            st.write("This app helps you gather information and brainstorm ideas for your website.")

            website_name = st.text_input("Website Name:")

        # Target audience selection (multiselect with default)
            target_audience_options = ["Everyone", "Students", "Professionals", "Tech Enthusiasts"]
            target_audience = st.multiselect("Who is your target audience?", target_audience_options, default=["Everyone"])

            website_goal = st.selectbox("What is the main goal of your website?",
                                    ["Inform", "Sell products/services", "Build a community"])

        # Sample content suggestions based on selections with informative error handling
            def generate_content_suggestions(target_audience, website_goal):
                content_suggestions = {
                    "Everyone": {
                        "Inform": ["Blog posts", "FAQs", "About Us page"],
                        "Sell products/services": ["Product descriptions", "Customer testimonials", "Pricing page"],
                        "Build a community": ["Forum", "Events page", "Member directory"],
                    },
                    "Students": {
                        "Inform": ["Course information", "Faculty profiles", "Campus resources"],
                        "Sell products/services": ["Student discounts", "Career services", "University merchandise"],
                        "Build a community": ["Student clubs", "Social events calendar", "Online forum"],
                    },
                    "Tech Enthusiasts": {  # Add suggestions for Tech Enthusiasts
                        "Inform": ["Tech tutorials", "Product reviews", "Industry news"],
                        "Sell products/services": ["Software subscriptions", "Hardware components", "Online courses"],
                        "Build a community": ["Tech forums", "Hackathons", "Meetup groups"],
                    },
                    "Professionals": {  # Add suggestions for Professionals
                        "Inform": ["Industry reports", "White papers", "Case studies"],
                        "Sell products/services": ["B2B solutions", "Pricing plans", "Client testimonials"],
                        "Build a community": ["Professional networking", "Industry events", "Online forums"],
                    }
                }

                # Handle empty or unsupported target audience gracefully
                if not target_audience:
                    return ["Please select a target audience."]
                elif target_audience[0] not in content_suggestions:
                    return [f"Content suggestions for '{target_audience[0]}' are not currently available. "
                            f"Please select another target audience or consider adding suggestions for '{target_audience[0]}' "
                            f"to the 'content_suggestions' dictionary."]  # Informative message

                selected_audience = target_audience.pop()
                return choices(content_suggestions[selected_audience][website_goal], k=3)  # Randomly pick 3 suggestions

            content_ideas = generate_content_suggestions(target_audience, website_goal)

            st.subheader("Website Content Inspiration:")
            st.write("Here are some content ideas based on your target audience and website goal:")
            for idea in content_ideas:
                st.write("- " + idea)

           
       
if page=="voiceassistance":

    
    import streamlit as st
    import speech_recognition as sr
    import pywhatkit

    def get_audio():
        recorder = sr.Recognizer()
        with sr.Microphone() as source:
            st.title(":blue Voice Assistant")
            st.markdown("---")  # Horizontal line
            st.info(":green This Voice Assitant can perform tasks based on your voice commands.")
            st.subheader(":orange Click on The Microphonr icon below, Speach and Watch it Response!")
            st.info("* For Example, You can ask it to search the Web Page or Play a Youtube video..")
            st.image("img.png", width=100)  # Add a microphone icon

            audio = recorder.listen(source)

        try:
            text = recorder.recognize_google(audio)
            st.write(f":orange [You said: {text}]")
            return text
        except sr.UnknownValueError:
            st.error("Sorry, I couldn't understand what you said.")
            return None
        except sr.RequestError:
            st.error("Sorry, I couldn't request results from the speech recognition service.")
            return None

    text = get_audio()

    if text:
        if "youtube" in text.lower():
            pywhatkit.playonyt(text)
        else:
            pywhatkit.search(text)




if page=="video":
    import streamlit as st
    from st_clickable_images import clickable_images
    import pandas as pd
    from pytube import YouTube
    import os
    import requests
    from time import sleep

    upload_endpoint = "https://api.assemblyai.com/v2/upload"
    transcript_endpoint = "https://api.assemblyai.com/v2/transcript"

    headers = {
        "authorization": "df100d0c58744221939e4b67757d4499",
        "content-type": "application/json"
    }

    @st.cache_data
    def save_audio(url):
        yt = YouTube(url)
        try:
            video = yt.streams.filter(only_audio=True).first()
            out_file = video.download()
        except:
            return None, None, None
        base, ext = os.path.splitext(out_file)
        file_name = base + '.mp3'
        os.rename(out_file, file_name)
        print(yt.title + " has been successfully downloaded.")
        print(file_name)
        return yt.title, file_name, yt.thumbnail_url

    @st.cache_data
    def upload_to_AssemblyAI(save_location):
        CHUNK_SIZE = 5242880
        print(save_location)

        def read_file(filename):
            with open(filename, 'rb') as _file:
                while True:
                    print("chunk uploaded")
                    data = _file.read(CHUNK_SIZE)
                    if not data:
                        break
                    yield data

        upload_response = requests.post(
            upload_endpoint,
            headers=headers, data=read_file(save_location)
        )
        print(upload_response.json())

        if "error" in upload_response.json():
            return None, upload_response.json()["error"]

        audio_url = upload_response.json()['upload_url']
        print('Uploaded to', audio_url)

        return audio_url, None

    @st.cache_data
    def start_analysis(audio_url):
        print(audio_url)

        ## Start transcription job of audio file
        data = {
            'audio_url': audio_url,
            'iab_categories': True,
            'content_safety': True,
            "summarization": True,
            "summary_model": "informative",
            "summary_type": "bullets"
        }

        transcript_response = requests.post(transcript_endpoint, json=data, headers=headers)
        print(transcript_response.json())

        if 'error' in transcript_response.json():
            return None, transcript_response.json()['error']

        transcript_id = transcript_response.json()['id']
        polling_endpoint = transcript_endpoint + "/" + transcript_id

        print("Transcribing at", polling_endpoint)
        return polling_endpoint, None

    @st.cache_data
    def get_analysis_results(polling_endpoint):

        status = 'submitted'

        while True:
            print(status)
            polling_response = requests.get(polling_endpoint, headers=headers)
            status = polling_response.json()['status']
            # st.write(polling_response.json())
            # st.write(status)

            if status == 'submitted' or status == 'processing' or status == 'queued':
                print('not ready yet')
                sleep(10)

            elif status == 'completed':
                print('creating transcript')

                return polling_response

                break
            else:
                print('error')
                return False
                break

    st.title("YouTube Content Analyzer")
    st.markdown("With this app you can audit a Youtube channel to see if you'd like to sponsor them. All you have to do is to pass a list of links to the videos of this channel and you will get a list of thumbnails. Once you select a video by clicking its thumbnail, you can view:")
    st.markdown("1. a summary of the video,") 
    st.markdown("2. the topics that are discussed in the video,") 
    st.markdown("3. whether there are any sensitive topics discussed in the video.")
    st.markdown("Make sure your links are in the format: https://www.youtube.com/watch?v=HfNnuQOHAaw and not https://youtu.be/HfNnuQOHAaw")

    default_bool = st.checkbox("Use a default file")

    if default_bool:
        file = open("E:\\Final_Project_Sem8\\links.txt")
    else:
        file = st.file_uploader("Upload a file that includes the links (.txt)")

    if file is not None:
        dataframe = pd.read_csv(file, header=None)
        dataframe.columns = ['urls']
        urls_list = dataframe['urls'].tolist()

        titles = []
        locations = []
        thumbnails = []

        for video_url in urls_list:
            # download audio
            video_title, save_location, video_thumbnail = save_audio(video_url)
            if video_title:
                titles.append(video_title)
                locations.append(save_location)
                thumbnails.append(video_thumbnail)

        selected_video = clickable_images(thumbnails,
        titles = titles,
        div_style={"height": "400px", "display": "flex", "justify-content": "center", "flex-wrap": "wrap", "overflow-y":"auto"},
        img_style={"margin": "5px", "height": "150px"}
        )

        st.markdown(f"Thumbnail #{selected_video} clicked" if selected_video > -1 else "No image clicked")

        if selected_video > -1:
            video_url = urls_list[selected_video]
            video_title = titles[selected_video]
            save_location = locations[selected_video]

            st.header(video_title)
            st.audio(save_location)

            # upload mp3 file to AssemblyAI
            audio_url, error = upload_to_AssemblyAI(save_location)
            
            if error:
                st.write(error)
            else:
                # start analysis of the file
                polling_endpoint, error = start_analysis(audio_url)

                if error:
                    st.write(error)
                else:
                    # receive the results
                    results = get_analysis_results(polling_endpoint)

                    summary = results.json()['summary']
                    topics = results.json()['iab_categories_result']['summary']
                    sensitive_topics = results.json()['content_safety_labels']['summary']

                    st.header("Summary of this video")
                    st.write(summary)

                    st.header("Sensitive content")
                    if sensitive_topics != {}:
                        st.subheader('🚨 Mention of the following sensitive topics detected.')
                        moderation_df = pd.DataFrame(sensitive_topics.items())
                        moderation_df.columns = ['topic','confidence']
                        st.dataframe(moderation_df, use_container_width=True)

                    else:
                        st.subheader('✅ All clear! No sensitive content detected.')

                    st.header("Topics discussed")
                    topics_df = pd.DataFrame(topics.items())
                    topics_df.columns = ['topic','confidence']
                    topics_df["topic"] = topics_df["topic"].str.split(">")
                    expanded_topics = topics_df.topic.apply(pd.Series).add_prefix('topic_level_')
                    topics_df = topics_df.join(expanded_topics).drop('topic', axis=1).sort_values(['confidence'], ascending=False).fillna('')

                    st.dataframe(topics_df)

    
if page=="Audio":
    import streamlit as st
    import assemblyai as aai

    def transcribe_audio(audio_url):
        config = aai.TranscriptionConfig(
        redact_pii=True,
        redact_pii_audio=True,
        redact_pii_policies=[
            aai.PIIRedactionPolicy.person_name,
            aai.PIIRedactionPolicy.phone_number,
        ],
        redact_pii_sub=aai.PIISubstitutionPolicy.hash,
    )
        transcript = aai.Transcriber().transcribe(audio_url, config)
        return transcript

    def main():
        st.title("Dynamic Audio Transcription and Redaction")

        # User input for audio file URL
        audio_url = st.text_input("Enter the URL of the MP3 audio file:")

        # Checkbox for redaction options
        #redact_pii = st.checkbox("Redact Personally Identifiable Information")
        redact_pii_audio = st.checkbox("Redact Audio")
    

        # Transcribe and redact audio
        if st.button("Transcribe and Redact"):
            if audio_url:
                transcript = transcribe_audio(audio_url)
                st.subheader("Transcript")
                st.write(transcript.text)
                if redact_pii_audio:
                    st.subheader("Redacted Audio")
                    st.audio(transcript.get_redacted_audio_url())

    if __name__  == '__main__':
        main()
